# encoding: utf-8

"""
Test suite for pptx.parts.slide module
"""

from __future__ import (
    absolute_import, division, print_function, unicode_literals
)

import pytest

from pptx.chart.data import ChartData
from pptx.enum.base import EnumValue
from pptx.opc.constants import CONTENT_TYPE as CT, RELATIONSHIP_TYPE as RT
from pptx.opc.packuri import PackURI
from pptx.oxml.slide import CT_Slide
from pptx.package import Package
from pptx.parts.chart import ChartPart
from pptx.parts.image import Image, ImagePart
from pptx.parts.presentation import PresentationPart
from pptx.parts.slide import (
    BaseSlidePart, SlideLayoutPart, SlideMasterPart, SlidePart
)
from pptx.slide import Slide, SlideLayout, SlideMaster

from ..unitutil.cxml import element
from ..unitutil.file import absjoin, test_file_dir
from ..unitutil.mock import (
    class_mock, initializer_mock, instance_mock, method_mock, property_mock
)


class DescribeBaseSlidePart(object):

    def it_knows_its_name(self, name_fixture):
        base_slide, expected_value = name_fixture
        assert base_slide.name == expected_value

    def it_can_get_a_related_image_by_rId(self, get_image_fixture):
        slide, rId, image_ = get_image_fixture
        assert slide.get_image(rId) is image_

    def it_can_add_an_image_part(self, image_part_fixture):
        slide, image_file, image_part_, rId_ = image_part_fixture

        image_part, rId = slide.get_or_add_image_part(image_file)

        slide._package.get_or_add_image_part.assert_called_once_with(
            image_file
        )
        slide.relate_to.assert_called_once_with(image_part_, RT.IMAGE)
        assert image_part is image_part_
        assert rId is rId_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def get_image_fixture(self, related_parts_prop_, image_part_, image_):
        slide = BaseSlidePart(None, None, None, None)
        rId = 'rId42'
        related_parts_prop_.return_value = {rId: image_part_}
        image_part_.image = image_
        return slide, rId, image_

    @pytest.fixture
    def image_part_fixture(
            self, partname_, package_, image_part_, relate_to_):
        slide = BaseSlidePart(partname_, None, None, package_)
        image_file, rId = 'foobar.png', 'rId6'
        package_.get_or_add_image_part.return_value = image_part_
        relate_to_.return_value = rId
        return slide, image_file, image_part_, rId

    @pytest.fixture
    def name_fixture(self):
        sld_cxml, expected_value = 'p:sld/p:cSld{name=Foobar}', 'Foobar'
        sld = element(sld_cxml)
        base_slide = BaseSlidePart(None, None, sld, None)
        return base_slide, expected_value

    # fixture components ---------------------------------------------

    @pytest.fixture
    def image_(self, request):
        return instance_mock(request, Image)

    @pytest.fixture
    def image_part_(self, request):
        return instance_mock(request, ImagePart)

    @pytest.fixture
    def package_(self, request):
        return instance_mock(request, Package)

    @pytest.fixture
    def partname_(self):
        return PackURI('/foo/bar.xml')

    @pytest.fixture
    def relate_to_(self, request):
        return method_mock(request, BaseSlidePart, 'relate_to')

    @pytest.fixture
    def related_parts_prop_(self, request):
        return property_mock(request, BaseSlidePart, 'related_parts')


class DescribeSlidePart(object):

    def it_knows_its_slide_id(self, slide_id_fixture):
        slide_part, presentation_part_, slide_id = slide_id_fixture
        _slide_id = slide_part.slide_id
        presentation_part_.slide_id.assert_called_once_with(slide_part)
        assert _slide_id is slide_id

    def it_can_add_a_chart_part(self, add_chart_part_fixture):
        slide_part, chart_type_, chart_data_ = add_chart_part_fixture[:3]
        ChartPart_, chart_part_, package_, rId = add_chart_part_fixture[3:]

        _rId = slide_part.add_chart_part(chart_type_, chart_data_)

        ChartPart_.new.assert_called_once_with(
            chart_type_, chart_data_, package_
        )
        slide_part.relate_to.assert_called_once_with(
            slide_part, chart_part_, RT.CHART
        )
        assert _rId is rId

    def it_can_create_a_new_slide_part(self, new_fixture):
        slide_layout_part_, partname, package_ = new_fixture[:3]
        SlidePart_init_, sld = new_fixture[3:]

        slide_part = SlidePart.new(partname, package_, slide_layout_part_)

        SlidePart_init_.assert_called_once_with(
            partname, CT.PML_SLIDE, sld, package_
        )
        slide_part.relate_to.assert_called_once_with(
            slide_part, slide_layout_part_, RT.SLIDE_LAYOUT
        )
        assert isinstance(slide_part, SlidePart)

    def it_provides_access_to_its_slide(self, slide_fixture):
        slide_part, Slide_, sld, slide_ = slide_fixture
        slide = slide_part.slide
        Slide_.assert_called_once_with(sld, slide_part)
        assert slide is slide_

    def it_provides_access_to_the_slide_layout(self, layout_fixture):
        slide_part, slide_layout_ = layout_fixture
        slide_layout = slide_part.slide_layout
        slide_part.part_related_by.assert_called_once_with(RT.SLIDE_LAYOUT)
        assert slide_layout is slide_layout_

    def it_knows_the_minimal_element_xml_for_a_slide(self):
        path = absjoin(test_file_dir, 'minimal_slide.xml')
        sld = CT_Slide.new()
        with open(path, 'r') as f:
            expected_xml = f.read()
        assert sld.xml == expected_xml

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def add_chart_part_fixture(
            self, package_, chart_type_, chart_data_, ChartPart_,
            chart_part_, relate_to_):
        slide_part = SlidePart(None, None, None, package_)
        ChartPart_.new.return_value = chart_part_
        relate_to_.return_value = rId = 'rId42'
        return (
            slide_part, chart_type_, chart_data_, ChartPart_, chart_part_,
            package_, rId
        )

    @pytest.fixture
    def layout_fixture(self, slide_layout_, part_related_by_):
        slide_part = SlidePart(None, None, None, None)
        part_related_by_.return_value.slide_layout = slide_layout_
        return slide_part, slide_layout_

    @pytest.fixture
    def new_fixture(
            self, slide_layout_part_, package_, SlidePart_init_, CT_Slide_,
            relate_to_):
        partname = PackURI('/foobar.xml')
        CT_Slide_.new.return_value = sld = element('c:sld')
        return slide_layout_part_, partname, package_, SlidePart_init_, sld

    @pytest.fixture
    def slide_fixture(self, Slide_, slide_):
        sld = element('p:sld')
        slide_part = SlidePart(None, None, sld, None)
        return slide_part, Slide_, sld, slide_

    @pytest.fixture
    def slide_id_fixture(self, package_, presentation_part_):
        slide_part = SlidePart(None, None, None, package_)
        slide_id = 256
        package_.presentation_part = presentation_part_
        presentation_part_.slide_id.return_value = slide_id
        return slide_part, presentation_part_, slide_id

    # fixture components -----------------------------------

    @pytest.fixture
    def ChartPart_(self, request, chart_part_):
        return class_mock(
            request, 'pptx.parts.slide.ChartPart', return_value=chart_part_
        )

    @pytest.fixture
    def chart_data_(self, request):
        return instance_mock(request, ChartData)

    @pytest.fixture
    def chart_part_(self, request):
        return instance_mock(request, ChartPart)

    @pytest.fixture
    def chart_type_(self, request):
        return instance_mock(request, EnumValue)

    @pytest.fixture
    def CT_Slide_(self, request):
        return class_mock(request, 'pptx.parts.slide.CT_Slide')

    @pytest.fixture
    def package_(self, request):
        return instance_mock(request, Package)

    @pytest.fixture
    def part_related_by_(self, request):
        return method_mock(request, SlidePart, 'part_related_by')

    @pytest.fixture
    def presentation_part_(self, request):
        return instance_mock(request, PresentationPart)

    @pytest.fixture
    def relate_to_(self, request):
        return method_mock(request, SlidePart, 'relate_to', autospec=True)

    @pytest.fixture
    def Slide_(self, request, slide_):
        return class_mock(
            request, 'pptx.parts.slide.Slide', return_value=slide_
        )

    @pytest.fixture
    def slide_(self, request):
        return instance_mock(request, Slide)

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def slide_layout_part_(self, request):
        return instance_mock(request, SlideLayoutPart)

    @pytest.fixture
    def SlidePart_init_(self, request):
        return initializer_mock(request, SlidePart)


class DescribeSlideLayoutPart(object):

    def it_provides_access_to_its_slide_master(self, master_fixture):
        slide_layout_part, part_related_by_, slide_master_ = master_fixture
        slide_master = slide_layout_part.slide_master
        part_related_by_.assert_called_once_with(
            slide_layout_part, RT.SLIDE_MASTER
        )
        assert slide_master is slide_master_

    def it_provides_access_to_its_slide_layout(self, layout_fixture):
        slide_layout_part, SlideLayout_ = layout_fixture[:2]
        sldLayout, slide_layout_ = layout_fixture[2:]
        slide_layout = slide_layout_part.slide_layout
        SlideLayout_.assert_called_once_with(sldLayout, slide_layout_part)
        assert slide_layout is slide_layout_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def layout_fixture(self, SlideLayout_, slide_layout_):
        sldLayout = element('p:sldLayout')
        slide_layout_part = SlideLayoutPart(None, None, sldLayout)
        return slide_layout_part, SlideLayout_, sldLayout, slide_layout_

    @pytest.fixture
    def master_fixture(
            self, part_related_by_, slide_master_, slide_master_part_):
        slide_layout_part = SlideLayoutPart(None, None, None, None)
        part_related_by_.return_value = slide_master_part_
        slide_master_part_.slide_master = slide_master_
        return slide_layout_part, part_related_by_, slide_master_

    # fixture components -----------------------------------

    @pytest.fixture
    def part_related_by_(self, request):
        return method_mock(
            request, SlideLayoutPart, 'part_related_by', autospec=True
        )

    @pytest.fixture
    def SlideLayout_(self, request, slide_layout_):
        return class_mock(
            request, 'pptx.parts.slide.SlideLayout',
            return_value=slide_layout_
        )

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def slide_master_(self, request):
        return instance_mock(request, SlideMaster)

    @pytest.fixture
    def slide_master_part_(self, request):
        return instance_mock(request, SlideMasterPart)


class DescribeSlideMasterPart(object):

    def it_provides_access_to_its_slide_master(self, master_fixture):
        slide_master_part, SlideMaster_, sldMaster, slide_master_ = (
            master_fixture
        )
        slide_master = slide_master_part.slide_master
        SlideMaster_.assert_called_once_with(sldMaster, slide_master_part)
        assert slide_master is slide_master_

    def it_provides_access_to_a_related_slide_layout(self, related_fixture):
        slide_master_part, rId, getitem_, slide_layout_ = related_fixture
        slide_layout = slide_master_part.related_slide_layout(rId)
        getitem_.assert_called_once_with(rId)
        assert slide_layout is slide_layout_

    # fixtures -------------------------------------------------------

    @pytest.fixture
    def master_fixture(self, SlideMaster_, slide_master_):
        sldMaster = element('p:sldMaster')
        slide_master_part = SlideMasterPart(None, None, sldMaster)
        return slide_master_part, SlideMaster_, sldMaster, slide_master_

    @pytest.fixture
    def related_fixture(self, slide_layout_, related_parts_prop_):
        slide_master_part = SlideMasterPart(None, None, None, None)
        rId = 'rId42'
        getitem_ = related_parts_prop_.return_value.__getitem__
        getitem_.return_value.slide_layout = slide_layout_
        return slide_master_part, rId, getitem_, slide_layout_

    # fixture components ---------------------------------------------

    @pytest.fixture
    def related_parts_prop_(self, request):
        return property_mock(request, SlideMasterPart, 'related_parts')

    @pytest.fixture
    def slide_layout_(self, request):
        return instance_mock(request, SlideLayout)

    @pytest.fixture
    def SlideMaster_(self, request, slide_master_):
        return class_mock(
            request, 'pptx.parts.slide.SlideMaster',
            return_value=slide_master_
        )

    @pytest.fixture
    def slide_master_(self, request):
        return instance_mock(request, SlideMaster)
